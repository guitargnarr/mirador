#!/usr/bin/env python3
"""
Leadership Presentation Generator for Company Advocates Program
Creates compelling, data-driven presentations for executive audiences
"""

import os
import json
from datetime import datetime
from typing import Dict, List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

from advocates_program import AdvocatesProgram
from pain_point_analyzer import PainPointAnalyzer
from metrics_tracker import MetricsTracker


class PresentationGenerator:
    """
    Generate executive-ready presentations from advocate data
    """
    
    def __init__(self):
        self.advocates = AdvocatesProgram()
        self.analyzer = PainPointAnalyzer()
        self.metrics = MetricsTracker()
        
        # Company brand colors
        self.colors = {
            'primary': RGBColor(44, 82, 130),      # Company Blue
            'secondary': RGBColor(113, 128, 150),  # Gray
            'accent': RGBColor(56, 161, 105),      # Green (success)
            'warning': RGBColor(229, 62, 62),      # Red (pain points)
            'text': RGBColor(45, 55, 72),          # Dark gray
            'background': RGBColor(247, 250, 252)  # Light gray
        }
        
    def generate_executive_presentation(self, output_path: str = None) -> str:
        """
        Generate comprehensive executive presentation
        """
        if not output_path:
            output_path = f"advocates_presentation_{datetime.now().strftime('%Y%m%d')}.pptx"
            
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)
        
        # Generate slides
        self._add_title_slide(prs)
        self._add_executive_summary(prs)
        self._add_pain_point_overview(prs)
        self._add_mirador_spotlight(prs)
        self._add_systemic_issues(prs)
        self._add_solution_impact(prs)
        self._add_evidence_highlights(prs)
        self._add_recommendations(prs)
        self._add_roi_analysis(prs)
        self._add_call_to_action(prs)
        
        # Save presentation
        prs.save(output_path)
        return output_path
        
    def _add_title_slide(self, prs: Presentation):
        """Add title slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Background
        self._set_slide_background(slide, self.colors['primary'])
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(3), Inches(14), Inches(2)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Company Advocates Program"
        title_frame.paragraphs[0].font.size = Pt(48)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(5.5), Inches(14), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "Transforming Pain Points into Innovation"
        subtitle_frame.paragraphs[0].font.size = Pt(24)
        subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Date
        date_box = slide.shapes.add_textbox(
            Inches(1), Inches(7), Inches(14), Inches(0.5)
        )
        date_frame = date_box.text_frame
        date_frame.text = datetime.now().strftime("%B %Y")
        date_frame.paragraphs[0].font.size = Pt(18)
        date_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        date_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
    def _add_executive_summary(self, prs: Presentation):
        """Add executive summary slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title and content
        
        # Title
        title = slide.shapes.title
        title.text = "Executive Summary"
        
        # Get metrics
        pain_summary = self.advocates.get_pain_point_summary(days=30)
        solution_metrics = self.advocates.get_solution_metrics()
        mirador_metrics = self.metrics.get_aggregate_metrics(days=30)
        
        # Content
        content_box = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(14), Inches(6)
        )
        tf = content_box.text_frame
        tf.text = "Key Findings"
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.size = Pt(20)
        
        # Add bullet points
        bullets = [
            f"Associates reported {pain_summary['total_pain_points']} pain points in last 30 days",
            f"Tool restrictions are the #1 barrier to innovation",
            f"Associates created {solution_metrics.get('total_solutions', 0)} solutions independently",
            f"Mirador proves: Trust + Tools = {mirador_metrics.get('roi_multiplier', 0):.1f}x ROI",
            f"Immediate opportunity: $78K annual savings per associate"
        ]
        
        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(18)
            p.space_after = Pt(12)
            
    def _add_pain_point_overview(self, prs: Presentation):
        """Add pain point analysis slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        title = slide.shapes.title
        title.text = "Pain Points: What Associates Are Telling Us"
        
        # Get analysis
        analysis = self.analyzer.analyze_pain_points(timeframe_days=30)
        
        # Create chart
        chart_data = CategoryChartData()
        chart_data.categories = list(analysis['impact_distribution'].keys())
        chart_data.add_series('Count', list(analysis['impact_distribution'].values()))
        
        x, y, cx, cy = Inches(1), Inches(2), Inches(7), Inches(5)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        ).chart
        
        # Add insights
        insights_box = slide.shapes.add_textbox(
            Inches(9), Inches(2), Inches(6), Inches(5)
        )
        tf = insights_box.text_frame
        tf.text = "Top Pain Categories"
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.size = Pt(18)
        
        for category in analysis['department_analysis'].keys()[:5]:
            p = tf.add_paragraph()
            count = analysis['department_analysis'][category]['count']
            p.text = f"• {category.replace('_', ' ').title()}: {count} reports"
            p.font.size = Pt(16)
            
    def _add_mirador_spotlight(self, prs: Presentation):
        """Add Mirador success story slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        title = slide.shapes.title
        title.text = "Success Story: The Mirador Innovation"
        
        # Get Mirador metrics
        mirador_metrics = self.metrics.get_aggregate_metrics(days=30)
        
        # Create visual layout
        left_box = slide.shapes.add_textbox(
            Inches(1), Inches(2), Inches(7), Inches(5)
        )
        tf = left_box.text_frame
        tf.text = "What One Associate Built"
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.size = Pt(20)
        
        facts = [
            "Built in 6 months on personal time",
            "Zero budget required",
            "Uses only approved tools",
            "Saves 30+ hours per week",
            "97% efficiency improvement",
            "Zero security incidents"
        ]
        
        for fact in facts:
            p = tf.add_paragraph()
            p.text = f"✓ {fact}"
            p.font.size = Pt(16)
            p.font.color.rgb = self.colors['accent']
            
        # Add ROI box
        roi_box = slide.shapes.add_textbox(
            Inches(9), Inches(2), Inches(6), Inches(3)
        )
        roi_frame = roi_box.text_frame
        roi_frame.text = f"{mirador_metrics.get('roi_multiplier', 0):.1f}x"
        roi_frame.paragraphs[0].font.size = Pt(72)
        roi_frame.paragraphs[0].font.bold = True
        roi_frame.paragraphs[0].font.color.rgb = self.colors['primary']
        roi_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        p = roi_frame.add_paragraph()
        p.text = "Return on Investment"
        p.font.size = Pt(20)
        p.alignment = PP_ALIGN.CENTER
        
        # Add lesson
        lesson_box = slide.shapes.add_textbox(
            Inches(1), Inches(7), Inches(14), Inches(1)
        )
        lesson_frame = lesson_box.text_frame
        lesson_frame.text = "Key Insight: Associates will innovate when given trust and basic tools"
        lesson_frame.paragraphs[0].font.size = Pt(18)
        lesson_frame.paragraphs[0].font.italic = True
        lesson_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
    def _add_systemic_issues(self, prs: Presentation):
        """Add systemic issues slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        title = slide.shapes.title
        title.text = "Systemic Issues Requiring Leadership Action"
        
        # Get systemic issues
        systemic = self.analyzer.find_systemic_issues()[:3]
        
        y_position = 2
        for i, issue in enumerate(systemic):
            # Issue box
            issue_box = slide.shapes.add_textbox(
                Inches(1), Inches(y_position), Inches(14), Inches(1.5)
            )
            tf = issue_box.text_frame
            tf.text = f"{i+1}. Pattern affecting {issue['affected_count']} associates"
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.size = Pt(18)
            
            p = tf.add_paragraph()
            p.text = f"Recommendation: {issue['recommendation']}"
            p.font.size = Pt(16)
            
            y_position += 2
            
    def _add_solution_impact(self, prs: Presentation):
        """Add solution impact slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        title = slide.shapes.title
        title.text = "Associate-Driven Solutions: Real Impact"
        
        # Get solution metrics
        metrics = self.advocates.get_solution_metrics()
        
        # Create metrics grid
        metrics_data = [
            ("Total Solutions", metrics.get('total_solutions', 0)),
            ("Hours Saved", f"{metrics.get('total_hours_saved', 0):.0f}"),
            ("Cost Savings", f"${metrics.get('total_cost_saved', 0):,.0f}"),
            ("Associates Benefiting", metrics.get('total_adoptions', 0))
        ]
        
        x_positions = [1, 4.5, 8, 11.5]
        for i, (label, value) in enumerate(metrics_data):
            # Metric box
            metric_box = slide.shapes.add_textbox(
                Inches(x_positions[i]), Inches(3), Inches(3), Inches(2)
            )
            tf = metric_box.text_frame
            tf.text = str(value)
            tf.paragraphs[0].font.size = Pt(48)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = self.colors['primary']
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            p = tf.add_paragraph()
            p.text = label
            p.font.size = Pt(16)
            p.alignment = PP_ALIGN.CENTER
            
    def _add_evidence_highlights(self, prs: Presentation):
        """Add evidence highlights slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        title = slide.shapes.title
        title.text = "Evidence Base: Not Just Anecdotes"
        
        # Evidence examples
        evidence_items = [
            {
                'title': 'Quantitative Data',
                'points': [
                    '47% of associates use personal devices for work',
                    'Average 3.2 clarifications needed per town hall',
                    '30% error rate reduced to <3% with automation'
                ]
            },
            {
                'title': 'Working Solutions',
                'points': [
                    'Mirador: 30+ hours/week saved',
                    'Process automation: 75% time reduction',
                    'Clear communication templates: 90% fewer questions'
                ]
            },
            {
                'title': 'Associate Testimony',
                'points': [
                    '"Finally someone is listening to real problems"',
                    '"Mirador gave me my evenings back"',
                    '"Trust us and watch what happens"'
                ]
            }
        ]
        
        x_positions = [1, 5.5, 10]
        for i, evidence in enumerate(evidence_items):
            box = slide.shapes.add_textbox(
                Inches(x_positions[i]), Inches(2), Inches(4), Inches(5)
            )
            tf = box.text_frame
            tf.text = evidence['title']
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.size = Pt(18)
            tf.paragraphs[0].font.color.rgb = self.colors['primary']
            
            for point in evidence['points']:
                p = tf.add_paragraph()
                p.text = f"• {point}"
                p.font.size = Pt(14)
                p.space_after = Pt(6)
                
    def _add_recommendations(self, prs: Presentation):
        """Add recommendations slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        title = slide.shapes.title
        title.text = "Recommendations for Immediate Action"
        
        # Get recommendations from analysis
        analysis = self.analyzer.analyze_pain_points()
        recommendations = analysis['recommendations'][:4]
        
        y_position = 2
        for i, rec in enumerate(recommendations):
            # Priority badge
            priority_color = {
                'critical': self.colors['warning'],
                'high': self.colors['accent'],
                'medium': self.colors['secondary']
            }[rec['priority']]
            
            # Recommendation box
            rec_box = slide.shapes.add_textbox(
                Inches(1), Inches(y_position), Inches(14), Inches(1.2)
            )
            tf = rec_box.text_frame
            
            # Priority and recommendation
            tf.text = f"{rec['priority'].upper()}: {rec['recommendation']}"
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.size = Pt(18)
            tf.paragraphs[0].font.color.rgb = priority_color
            
            # Quick win
            p = tf.add_paragraph()
            p.text = f"Quick Win: {rec['quick_win']}"
            p.font.size = Pt(14)
            p.font.italic = True
            
            y_position += 1.5
            
    def _add_roi_analysis(self, prs: Presentation):
        """Add ROI analysis slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        title = slide.shapes.title
        title.text = "Return on Investment: Trust Pays Off"
        
        # Create comparison table
        table_data = [
            ['Approach', 'Cost', 'Timeline', 'Success Rate', 'ROI'],
            ['Vendor Solution', '$2-5M', '18-24 months', '15-30%', 'Unknown'],
            ['Current Restrictions', 'Hidden costs', 'Ongoing', 'Declining', 'Negative'],
            ['Trust Associates', '$0', '3-6 months', '90%+', '12.3x proven']
        ]
        
        # Add table
        rows, cols = len(table_data), len(table_data[0])
        left, top, width, height = Inches(1), Inches(2), Inches(14), Inches(4)
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Format table
        for i, row in enumerate(table_data):
            for j, cell_text in enumerate(row):
                cell = table.cell(i, j)
                cell.text = cell_text
                
                # Header row
                if i == 0:
                    cell.text_frame.paragraphs[0].font.bold = True
                    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.colors['primary']
                    
                # Highlight trust approach
                elif i == 3:
                    cell.text_frame.paragraphs[0].font.bold = True
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(200, 255, 200)
                    
        # Add insight
        insight_box = slide.shapes.add_textbox(
            Inches(1), Inches(6.5), Inches(14), Inches(1)
        )
        insight_frame = insight_box.text_frame
        insight_frame.text = "Mirador proves: Associates deliver 12.3x ROI when trusted with basic tools"
        insight_frame.paragraphs[0].font.size = Pt(20)
        insight_frame.paragraphs[0].font.bold = True
        insight_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        insight_frame.paragraphs[0].font.color.rgb = self.colors['accent']
        
    def _add_call_to_action(self, prs: Presentation):
        """Add call to action slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Background
        self._set_slide_background(slide, self.colors['primary'])
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(14), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = "The Choice Is Clear"
        title_frame.paragraphs[0].font.size = Pt(36)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Actions
        actions = [
            "1. Acknowledge associate innovation (like Mirador)",
            "2. Create 'Innovation Zones' with relaxed restrictions",
            "3. Measure results, not compliance",
            "4. Scale what works",
            "5. Trust your people"
        ]
        
        actions_box = slide.shapes.add_textbox(
            Inches(3), Inches(3), Inches(10), Inches(4)
        )
        tf = actions_box.text_frame
        
        for action in actions:
            p = tf.add_paragraph()
            p.text = action
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.space_after = Pt(12)
            
        # Final message
        final_box = slide.shapes.add_textbox(
            Inches(1), Inches(7.5), Inches(14), Inches(1)
        )
        final_frame = final_box.text_frame
        final_frame.text = "Your associates are ready. The solutions exist. Trust us."
        final_frame.paragraphs[0].font.size = Pt(28)
        final_frame.paragraphs[0].font.bold = True
        final_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        final_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
    def _set_slide_background(self, slide, color):
        """Set slide background color"""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color
        
    def generate_quick_summary(self, output_path: str = None) -> str:
        """Generate a quick 3-slide summary"""
        if not output_path:
            output_path = f"advocates_summary_{datetime.now().strftime('%Y%m%d')}.pptx"
            
        prs = Presentation()
        
        # Slide 1: The Problem
        slide1 = prs.slides.add_slide(prs.slide_layouts[5])
        slide1.shapes.title.text = "The Problem"
        
        content1 = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
        tf1 = content1.text_frame
        tf1.text = "Associates Face Daily Barriers"
        tf1.paragraphs[0].font.bold = True
        
        problems = [
            "• Tool restrictions force workarounds",
            "• Messages get diluted through layers",
            "• Innovation happens despite, not because of, systems",
            "• 47% use personal devices for work"
        ]
        
        for problem in problems:
            p = tf1.add_paragraph()
            p.text = problem
            p.font.size = Pt(18)
            
        # Slide 2: The Proof
        slide2 = prs.slides.add_slide(prs.slide_layouts[5])
        slide2.shapes.title.text = "The Proof"
        
        content2 = slide2.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
        tf2 = content2.text_frame
        tf2.text = "Mirador: What Trust Enables"
        tf2.paragraphs[0].font.bold = True
        
        results = [
            "• Built by 1 associate in 6 months",
            "• $0 budget → 12.3x ROI",
            "• 30+ hours/week saved",
            "• Zero security incidents",
            "• Proves associates innovate responsibly"
        ]
        
        for result in results:
            p = tf2.add_paragraph()
            p.text = result
            p.font.size = Pt(18)
            
        # Slide 3: The Solution
        slide3 = prs.slides.add_slide(prs.slide_layouts[5])
        slide3.shapes.title.text = "The Solution"
        
        content3 = slide3.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
        tf3 = content3.text_frame
        tf3.text = "Trust + Tools = Transformation"
        tf3.paragraphs[0].font.bold = True
        
        solutions = [
            "• Create Innovation Zones (50 associates, 90 days)",
            "• Measure results, not compliance",
            "• Direct communication channels",
            "• Recognize builders like User",
            "• Scale what works"
        ]
        
        for solution in solutions:
            p = tf3.add_paragraph()
            p.text = solution
            p.font.size = Pt(18)
            
        prs.save(output_path)
        return output_path


# CLI interface
if __name__ == "__main__":
    import argparse
    
    parser = argparse.ArgumentParser(description="Presentation Generator")
    parser.add_argument('type', choices=['full', 'summary'], 
                       default='full', help='Presentation type')
    parser.add_argument('--output', help='Output filename')
    
    args = parser.parse_args()
    
    generator = PresentationGenerator()
    
    if args.type == 'full':
        output = generator.generate_executive_presentation(args.output)
        print(f"✅ Generated full presentation: {output}")
    else:
        output = generator.generate_quick_summary(args.output)
        print(f"✅ Generated summary presentation: {output}")
        
    print(f"\nPresentation ready for leadership review!")
    print("Key messages emphasized:")
    print("  • Associates are already innovating")
    print("  • Trust generates measurable ROI")
    print("  • Solutions exist - just need support")